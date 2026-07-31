#!/usr/bin/env python3
"""Generate the optional My Voice (Cox's Bazar) session slides — one idea per slide, big and minimal,
for the minority of sites with a screen. Fully offline; editable.

Run from anywhere:  python3 scripts/downloads/generate-slides.py  (needs: pip install python-pptx)
Override the output dir:  OUT_DIR=/tmp/pack python3 scripts/downloads/generate-slides.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
OUT_DIR = os.environ.get("OUT_DIR") or os.path.join(ROOT, "public", "downloads")
OUT = os.path.join(OUT_DIR, "my-voice-slides.pptx")
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
PLUM = RGBColor(0x7A, 0x3B, 0x69)
GREY = RGBColor(0x5A, 0x64, 0x73)
CREAM = RGBColor(0xF7, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, text, *, size, color, top, height, bold=True, align=PP_ALIGN.CENTER, left=Inches(0.8), width=None):
    width = width or (SW - Inches(1.6))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def title_slide(kicker, title, sub, accent=NAVY):
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    textbox(s, kicker, size=22, color=PLUM, top=Inches(2.2), height=Inches(0.8))
    textbox(s, title, size=60, color=accent, top=Inches(2.9), height=Inches(1.8))
    if sub:
        textbox(s, sub, size=26, color=GREY, top=Inches(4.7), height=Inches(1.2), bold=False)


def idea_slide(title, big, note=None, accent=NAVY):
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    # accent bar
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.35))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    textbox(s, title, size=30, color=PLUM, top=Inches(0.7), height=Inches(1.0))
    textbox(s, big, size=54, color=accent, top=Inches(2.2), height=Inches(2.6))
    if note:
        textbox(s, note, size=24, color=GREY, top=Inches(5.2), height=Inches(1.4), bold=False)


title_slide("Learning Bridge+  ·  Cox's Bazar", "My Voice", "Session slides — say who you are, in English")
idea_slide("Welcome", "Hello.\nThis is our space.", "Everyone here has a voice. We will find the English to say who we are.")
idea_slide("Our classroom words", "hello   ·   thank you\nplease   ·   again", "Say them with me. Use them every day.")
idea_slide("Listening: hello and names", "Listen.\nWhose name do you hear?", "Point. Stand. Show me — no writing yet.")
idea_slide("The alphabet and your name", "What is the first\nsound of your name?", "Fatima — fff. Say your name, slowly.")
idea_slide("Key sounds", "a  o  u\ns  t  l  m  n\np  t  k   ·   b  d  g", "Say the sound, not the letter name.")
idea_slide("Sound and letter practice", "Listen… which sound?\nBuild your word.", "Games. Play them again and again.")
idea_slide("Words about me", "my name · my family\nwhere I am from · what I like", "Draw it. Say it. You choose what to share.")
idea_slide("Writing my name", "I can write\nmy name.", "Trace it. Again. It is yours.")
idea_slide("I am…", "I am ______\nI am from ______\nI like ______", "Say a true sentence about you.", accent=PLUM)
idea_slide("Sentence practice", "I  am\nyou / we / they  are\nhe / she / it  is", "am, is, are — say it, build it.", accent=PLUM)
idea_slide("Meeting people", "Hello. My name is ______.\nWhat is your name?", "Greet someone. Introduce yourself.", accent=PLUM)
idea_slide("Make your card", "My Name,\nMy Voice", "Your name, big. Draw and write what shows who you are.")
idea_slide("Share", "This is me.\nThis is my voice.", "Show your card. We listen. We clap.")
title_slide("Well done", "Look how your\nvoice has grown", "From no English to saying who you are.", accent=PLUM)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("wrote", OUT, f"({len(prs.slides._sldIdLst)} slides)")
